import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

df = pd.read_csv('ufc-fighters-statistics.csv')

df['stance'] = df['stance'].fillna('Unspecified')

# calculate the day of birth
df['date_of_birth'] = pd.to_datetime(df['date_of_birth'], format='%Y-%m-%d', errors='coerce')
df['age'] = df['date_of_birth'].apply(lambda x: datetime.now().year - x.year if pd.notnull(x) else None)

# Handle the placeholder date by setting the age to None for those entries
df['age'] = df['age'].replace(datetime.now().year - 1900, None)

# Create a PowerPoint presentation object
prs = Presentation()


def save_plot_to_png(fig, slide_title):
    img_path = f"{slide_title}.png"
    fig.savefig(img_path)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_title
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8.5), height=Inches(4.75))

# Stance Distribution
fig = plt.figure()
sns.countplot(data=df, x='stance', color='green')
plt.title('Stance Distribution')
save_plot_to_png(fig, 'Stance Distribution')
plt.close(fig)

# Stance vs. Strikes Landed
fig = plt.figure()
sns.boxplot(data=df, x='stance', y='significant_strikes_landed_per_minute', color='green')
plt.title('Stance vs. Strikes Landed')
save_plot_to_png(fig, 'Stance vs. Strikes Landed')
plt.close(fig)

# Distribution of Wins, Losses, and Draws
fig, axs = plt.subplots(1, 3, figsize=(18, 5))
sns.histplot(df['wins'], bins=20, kde=True, ax=axs[0], color='green').set(title='Wins Distribution')
sns.histplot(df['losses'], bins=20, kde=True, ax=axs[1], color='red').set(title='Losses Distribution')
sns.histplot(df['draws'], bins=20, kde=True, ax=axs[2], color='blue').set(title='Draws Distribution')
save_plot_to_png(fig, 'Wins, Losses, and Draws Distribution')
plt.close(fig)

# Height vs. Weight
fig = plt.figure()
sns.scatterplot(data=df, x='height_cm', y='weight_in_kg', hue='stance', palette='muted')
plt.title('Height vs. Weight')
save_plot_to_png(fig, 'Height vs. Weight')
plt.close(fig)

# Age Distribution
fig = plt.figure()
sns.histplot(df['age'], bins=20, kde=True, color='purple')
plt.title('Age Distribution')
save_plot_to_png(fig, 'Age Distribution')
plt.close(fig)

# Pairplot for striking metrics

fig = sns.pairplot(df, vars=['significant_strikes_landed_per_minute', 'significant_striking_accuracy', 'significant_strikes_absorbed_per_minute', 'significant_strike_defence'], hue='stance', palette='bright')
fig.fig.suptitle('Striking Metrics Pairplot', y=1.02)
save_plot_to_png(fig.fig, 'Striking Metrics Pairplot')
plt.close(fig.fig)

# Takedown metrics
fig, axs = plt.subplots(2, 2, figsize=(15, 10))
sns.boxplot(data=df, x='stance', y='average_takedowns_landed_per_15_minutes', ax=axs[0, 0])
axs[0, 0].set_title('Average Takedowns per 15 Minutes')
sns.boxplot(data=df, x='stance', y='takedown_accuracy', ax=axs[0, 1])
axs[0, 1].set_title('Takedown Accuracy')
sns.boxplot(data=df, x='stance', y='takedown_defense', ax=axs[1, 0])
axs[1, 0].set_title('Takedown Defense')
sns.boxplot(data=df, x='stance', y='average_submissions_attempted_per_15_minutes', ax=axs[1, 1])
axs[1, 1].set_title('Average Submissions per 15 Minutes')
save_plot_to_png(fig, 'Takedown Metrics')
plt.close(fig)

# Win Rate vs. Significant Striking Accuracy
fig = plt.figure()
df['total_fights'] = df['wins'] + df['losses'] + df['draws']
df['win_rate'] = df['wins'] / df['total_fights']
sns.scatterplot(data=df, x='significant_striking_accuracy', y='win_rate', hue='stance', palette='viridis')
plt.title('Win Rate vs. Significant Striking Accuracy')
save_plot_to_png(fig, 'Win Rate vs. Significant Striking Accuracy')
plt.close(fig)

# Reach and Weight Distribution by Stance
fig, axs = plt.subplots(1, 2, figsize=(15, 5))
sns.violinplot(data=df, x='stance', y='reach_in_cm', ax=axs[0])
axs[0].set_title('Reach Distribution by Stance')
sns.violinplot(data=df, x='stance', y='weight_in_kg', ax=axs[1])
axs[1].set_title('Weight Distribution by Stance')
save_plot_to_png(fig, 'Reach and Weight Distribution by Stance')
plt.close(fig)

prs.save('UFC_Fighter_Statistics_Visualizations.pptx')